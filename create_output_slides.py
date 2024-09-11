import os
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns

# Path to the research folder
image_folder = 'synthetic/'

# Path to the output folder
map_folder = 'output_synthetic/'

# Path to the metrics.csv file
metrics_file = 'output_synthetic\metrics.csv'

# Path to the predictions.csv file
predictions_file = 'output_synthetic\predictions.csv'

# Load metrics data from the CSV file
metrics_data = pd.read_csv(metrics_file)

# Load predictions data from the CSV file
predictions_data = pd.read_csv(predictions_file)

# Create compiled DataFrame with all data
##5% diameter
ID_AL = ['0-3ATYY4','0-6NBE7Z','0-P5E5IT','0-P9UB1H','0-Q020H5','0-Z58QNY']
##2% diameter
ID_AH = ['0-6JNS7M','0-CBTJTK','0-EIT0RD','0-P2Z3RZ','0-QI7FRK','0-UKVR3B']
##no cavities
ID_none = ['11 78K_2','11 39K','1 TEM-BF 94K defocus 0 um']
#classify fatigue tests by ID
##5% areal density
ID_FL = ['0-3ATYY4','0-6JNS7M','0-6NBE7Z','0-CBTJTK','0-P5E5IT','0-UKVR3B']
##25% areal density
ID_FH = ['0-EIT0RD','0-P2Z3RZ','0-P9UB1H','0-Q020H5','0-QI7FRK','0-Z58QNY']
#classify distraction tests by ID
##easy background
ID_DE = ['0-P2Z3RZ','0-P5E5IT','0-P9UB1H','0-UKVR3B','11 78K_2']
##medium background
ID_DM = ['0-3ATYY4','0-6JNS7M','0-EIT0RD','0-Q020H5','11 39K']
##difficult background
ID_DD = ['0-6NBE7Z','0-CBTJTK','0-QI7FRK','0-Z58QNY','1 TEM-BF 94K defocus 0 um']

# Create empty DataFrame
compiled_data = pd.DataFrame(columns=['Image', 'Diameter', 'Density', 'Qualitative Difficulty', 'Symmetry', 'Clutter', 'Prediction'])

# Create a new PowerPoint presentation
presentation = Presentation()


# Iterate over the images in the research folder
for filename in os.listdir(image_folder):
    if filename.endswith('.jpg') or filename.endswith('.png'):
        research_image_path = os.path.join(image_folder, filename)
        output_image_path = os.path.join(map_folder, filename)
        
        # Get the associated value from metrics data
        value = metrics_data.loc[metrics_data['file'] == filename].values[0]
        image = value[0]
        symm = round(value[1],2)
        clutter = value[2]

        # Get the associated prediction from predictions data
        prediction = predictions_data.loc[predictions_data['Images'] == filename[:-4]].values[0]
        ID = prediction[0][6:]
        back_ID = prediction[0]
        pred = prediction[1]
        
        # Create a new slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        
        # Add the research image to the slide
        slide.shapes.add_picture(research_image_path, Inches(4.35), Inches(0), width=Inches(5.65), height=Inches(3.77))
        
        # Add the output image to the slide
        slide.shapes.add_picture(output_image_path, Inches(4.35), Inches(3.77), width=Inches(5.65), height=Inches(3.77))
        
        # Add the value and prediction as text to the slide
        text_box = slide.shapes.add_textbox(Inches(0.23), Inches(3.08), Inches(3.73), Inches(1))
        text_box.text_frame.text = f"Image: {image}\nSymmetry: {symm}\nClutter: {clutter}\nPrediction: {pred}"

        # Add to compiled DataFrame
        if ID in ID_AL:
            Size = '5%'
        elif ID in ID_AH:
            Size = '2%'
        elif back_ID in ID_none:
            Size = '0%'
        if ID in ID_FL:
            Density = '5%'
        elif ID in ID_FH:
            Density = '25%'
        elif back_ID in ID_none:
            Density = '0%'
        if (ID in ID_DE) or (back_ID in ID_DE):
            Difficulty = 'Easy'
        elif (ID in ID_DM) or (back_ID in ID_DM):
            Difficulty = 'Moderate'
        elif (ID in ID_DD) or (back_ID in ID_DD):
            Difficulty = 'Difficult'
        data = {'Image': [ID], 'Diameter': [Size], 'Density': [Density], 'Qualitative Difficulty': [Difficulty], 'Symmetry': [symm], 'Clutter': [clutter], 'Prediction': [pred]}
        if compiled_data.empty:
            compiled_data = pd.DataFrame(data)
        else:
            compiled_data = pd.concat([compiled_data, pd.DataFrame(data)], ignore_index=True)

# Sort compiled data by difficulty
compiled_data['Qualitative Difficulty'] = pd.Categorical(compiled_data['Qualitative Difficulty'], ['Easy','Moderate','Difficult'])
compiled_data = compiled_data.sort_values('Qualitative Difficulty')

# Create scatter plots for the compiled data
sns.scatterplot(data=compiled_data,x='Qualitative Difficulty',y='Clutter',style='Density',hue='Diameter',palette='tab10',alpha=0.6)
plt.savefig('figures\difficulty_clutter.png',dpi=1000)
#plt.show()
plt.close()

sns.scatterplot(data=compiled_data,x='Qualitative Difficulty',y='Symmetry',style='Density',hue='Diameter',palette='tab10',alpha=0.6)
plt.savefig('figures\difficulty_symmetry.png',dpi=1000)
#plt.show()
plt.close()

sns.scatterplot(data=compiled_data,x='Qualitative Difficulty',y='Prediction',style='Density',hue='Diameter',palette='tab10',alpha=0.6)
plt.savefig('figures\difficulty_prediction.png',dpi=1000)
#plt.show()
plt.close()

# Add scatter plots to the PowerPoint presentation
slide = presentation.slides.add_slide(presentation.slide_layouts[6])
slide.shapes.add_picture('figures\difficulty_clutter.png', Inches(0.5), Inches(0.5), width=Inches(9), height=Inches(6))

slide = presentation.slides.add_slide(presentation.slide_layouts[6])
slide.shapes.add_picture('figures\difficulty_symmetry.png', Inches(0.5), Inches(0.5), width=Inches(9), height=Inches(6))

slide = presentation.slides.add_slide(presentation.slide_layouts[6])
slide.shapes.add_picture('figures\difficulty_prediction.png', Inches(0.5), Inches(0.5), width=Inches(9), height=Inches(6))

# Add table to the PowerPoint presentation
slide = presentation.slides.add_slide(presentation.slide_layouts[6])
x, y, cx, cy = Inches(0.12), Inches(0.17), Inches(9.76), Inches(7.16)
shape = slide.shapes.add_table(16, 7, x, y, cx, cy)
table = shape.table

for i, row in compiled_data.iterrows():
    for j, column in enumerate(compiled_data.columns):
        if i == 0:
            table.cell(i, j).text = column
        table.cell(i+1, j).text = str(row[column])

# Save the PowerPoint presentation
presentation.save('output_synthetic\summary.pptx')