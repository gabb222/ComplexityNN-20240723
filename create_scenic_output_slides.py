import os
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import seaborn as sns

# Path to the research folder
image_folder = 'VISCHEMA-data/VISC-C/scenes/'

# Path to the scenic output folder
rgb_folder = 'output_scenic_rgb/'

# Path to the greyscale scenic output folder
grey_folder = 'output_scenic_grey/'

# Path to the metrics.csv file
metrics_file = 'metrics.csv'

# Path to the predictions.csv file
predictions_file = 'predictions.csv'

# Load metrics data from the CSV file
rgb_metrics_data = pd.read_csv(rgb_folder+metrics_file)
grey_metrics_data = pd.read_csv(grey_folder+metrics_file)

# Load predictions data from the CSV file
rgb_predictions_data = pd.read_csv(rgb_folder+predictions_file)
grey_predictions_data = pd.read_csv(grey_folder+predictions_file)

compiled_data = pd.DataFrame(columns=['Image', 'Color', 'Symmetry', 'Clutter', 'Prediction'])
compared_data = pd.DataFrame(columns=['Image', 'Symmetry Difference', 'Clutter Difference', 'Prediction Difference'])

# Iterate over the images in the research folder
for filename in os.listdir(image_folder):
    if filename.endswith('.jpg') or filename.endswith('.png'):

        # RGB
        # Get the associated value from metrics data
        value = rgb_metrics_data.loc[rgb_metrics_data['file'] == filename].values[0]
        image = value[0]
        symm1 = round(value[1],2)
        clutter1 = value[2]

        # Get the associated prediction from predictions data
        prediction = rgb_predictions_data.loc[rgb_predictions_data['Images'] == filename[:-4]].values[0]
        pred1 = prediction[1]

        data = {'Image': [image], 'Type': ['rgb'], 'Symmetry': [symm1], 'Clutter': [clutter1], 'Prediction': [pred1]}
        if compiled_data.empty:
            compiled_data = pd.DataFrame(data)
        else:
            compiled_data = pd.concat([compiled_data, pd.DataFrame(data)], ignore_index=True)

        # Greyscale
        # Get the associated value from metrics data
        value = grey_metrics_data.loc[grey_metrics_data['file'] == filename].values[0]
        image = value[0]
        symm2 = round(value[1],2)
        clutter2 = value[2]

        # Get the associated prediction from predictions data
        prediction = grey_predictions_data.loc[grey_predictions_data['Images'] == filename[:-4]].values[0]
        pred2 = prediction[1]

        data = {'Image': [image], 'Type': ['grey'], 'Symmetry': [symm2], 'Clutter': [clutter2], 'Prediction': [pred2]}
        if compiled_data.empty:
            compiled_data = pd.DataFrame(data)
        else:
            compiled_data = pd.concat([compiled_data, pd.DataFrame(data)], ignore_index=True)

        # Compare RGB and Greyscale
        comp_data = {'Image': [image], 'Symmetry Difference': [symm1-symm2], 'Clutter Difference': [clutter1-clutter2], 'Prediction Difference': [pred1-pred2],\
                    'RGB Symmetry': [symm1], 'Grey Symmetry': [symm2], 'RGB Clutter': [clutter1], 'Grey Clutter': [clutter2], 'RGB Prediction': [pred1], 'Grey Prediction': [pred2]}
        if compared_data.empty:
            compared_data = pd.DataFrame(comp_data)
        else:
            compared_data = pd.concat([compared_data, pd.DataFrame(comp_data)], ignore_index=True)



# Create scatterplot with seaborn
sns.scatterplot(data=compiled_data, x='Clutter', y='Prediction', hue='Type',alpha=0.6)
plt.savefig('figures/scenic_clutter_prediction_scatterplot.png')
plt.close()

sns.scatterplot(data=compiled_data, x='Symmetry', y='Prediction', hue='Type', alpha=0.6)
plt.savefig('figures/scenic_symmetry_prediction_scatterplot.png')
plt.close()

sns.scatterplot(data=compiled_data, x='Symmetry', y='Clutter', hue='Type', alpha=0.6)
plt.savefig('figures/scenic_symmetry_clutter_scatterplot.png', bbox_inches='tight')
plt.close()

# Create histogram with seaborn
fig, axs = plt.subplots(1, 3, figsize=(18, 5))

sns.violinplot(data=compared_data, y='Clutter Difference', orient='v', ax=axs[0], inner='quartiles', alpha=0.6)
sns.violinplot(data=compared_data, y='Symmetry Difference', orient='v', ax=axs[1], inner='quartiles', alpha=0.6)
sns.violinplot(data=compared_data, y='Prediction Difference', orient='v', ax=axs[2], inner='quartiles', alpha=0.6)

plt.savefig('figures/scenic_difference_violinplot.png', bbox_inches='tight')
plt.close()

# Create parity plot with seaborn of rgb predictions vs greyscale predictions
fig, axs = plt.subplots(1, 3, figsize=(18, 5))

sns.scatterplot(data=compared_data, x='RGB Clutter', y='Grey Clutter', ax=axs[0], alpha=0.6)
sns.scatterplot(data=compared_data, x='RGB Symmetry', y='Grey Symmetry', ax=axs[1], alpha=0.6)
sns.scatterplot(data=compared_data, x='RGB Prediction', y='Grey Prediction', ax=axs[2], alpha=0.6)

plt.savefig('figures/scenic_parity_plot.png', bbox_inches='tight')
plt.close()

# Create histogram with seaborn
fig, axs = plt.subplots(1, 3, figsize=(18, 5))

sns.violinplot(data=compiled_data, y='Clutter', orient='v', ax=axs[0], inner='quartiles', split=True, hue='Type', alpha=0.6)
sns.violinplot(data=compiled_data, y='Symmetry', orient='v', ax=axs[1], inner='quartiles', split=True, hue='Type', alpha=0.6)
sns.violinplot(data=compiled_data, y='Prediction', orient='v', ax=axs[2], inner='quartiles', split=True, hue='Type', alpha=0.6)

plt.savefig('figures/scenic_compare_violinplot.png', bbox_inches='tight')
plt.close()

# Create a new PowerPoint presentation
presentation = Presentation()

# Add scatter plots to the PowerPoint presentation
slide = presentation.slides.add_slide(presentation.slide_layouts[6])
slide.shapes.add_picture('figures\scenic_compare_violinplot.png', Inches(0.17), Inches(2.28), width=Inches(9.67), height=Inches(3.11))

slide = presentation.slides.add_slide(presentation.slide_layouts[6])
slide.shapes.add_picture('figures\scenic_parity_plot.png', Inches(0.17), Inches(2.28), width=Inches(9.67), height=Inches(3.11))

slide = presentation.slides.add_slide(presentation.slide_layouts[6])
slide.shapes.add_picture('figures\scenic_difference_violinplot.png', Inches(0.17), Inches(2.28), width=Inches(9.67), height=Inches(3.11))

slide = presentation.slides.add_slide(presentation.slide_layouts[6])
slide.shapes.add_picture('figures/scenic_symmetry_prediction_scatterplot.png', Inches(0), Inches(3.75), width=Inches(5.4), height=Inches(3.75))
slide.shapes.add_picture('figures/scenic_clutter_prediction_scatterplot.png', Inches(0), Inches(0), width=Inches(5.4), height=Inches(4.19))
slide.shapes.add_picture('figures/scenic_symmetry_clutter_scatterplot.png', Inches(4.95), Inches(2.19), width=Inches(5.05), height=Inches(3.59))

# Save the PowerPoint presentation
presentation.save('output_scenic_grey\summary.pptx')